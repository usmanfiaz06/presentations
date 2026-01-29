#!/usr/bin/env python3
"""
SERA 2026 Technical Proposal - PowerPoint Generator
Creates a PPT presentation from the HTML content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
DARK_BLUE = RGBColor(0x1a, 0x36, 0x5d)
LIGHT_BLUE = RGBColor(0x31, 0x82, 0xce)
GREEN = RGBColor(0x48, 0xbb, 0x78)
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xf7, 0xfa, 0xfc)

# Slide dimensions (16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def add_top_border(slide):
    """Add the top border design"""
    # Green bar on right
    green_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SLIDE_WIDTH - Inches(3.5), Inches(0),
        Inches(3.5), Inches(0.15)
    )
    green_bar.fill.solid()
    green_bar.fill.fore_color.rgb = GREEN
    green_bar.line.fill.background()

    # Dark blue bar below
    blue_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.2),
        SLIDE_WIDTH - Inches(1), Inches(0.08)
    )
    blue_bar.fill.solid()
    blue_bar.fill.fore_color.rgb = DARK_BLUE
    blue_bar.line.fill.background()

def add_bottom_border(slide, page_num):
    """Add the bottom border design with page number"""
    # Green line
    green_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), SLIDE_HEIGHT - Inches(0.55),
        SLIDE_WIDTH - Inches(1.6), Inches(0.05)
    )
    green_line.fill.solid()
    green_line.fill.fore_color.rgb = GREEN
    green_line.line.fill.background()

    # Dark blue line
    blue_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.8), SLIDE_HEIGHT - Inches(0.4),
        SLIDE_WIDTH - Inches(1.6), Inches(0.05)
    )
    blue_line.fill.solid()
    blue_line.fill.fore_color.rgb = DARK_BLUE
    blue_line.line.fill.background()

    # Page number
    page_box = slide.shapes.add_textbox(
        Inches(0.8), SLIDE_HEIGHT - Inches(0.5),
        Inches(0.6), Inches(0.4)
    )
    tf = page_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = f"{page_num:02d}"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE

def set_rtl_arabic_text(shape, text, font_size=18, bold=False, color=DARK_BLUE, align='right'):
    """Set Arabic text with RTL direction"""
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Arial'
    if align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    elif align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'left':
        p.alignment = PP_ALIGN.LEFT

def add_info_card(slide, left, top, width, height, title, value, bg_color=DARK_BLUE):
    """Add an info card with icon placeholder"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.4),
        width - Inches(0.2), Inches(0.4)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Value
    value_box = slide.shapes.add_textbox(
        left + Inches(0.1), top + Inches(0.7),
        width - Inches(0.2), Inches(0.5)
    )
    tf = value_box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_detail_box(slide, left, top, width, height, title, content, bg_color=LIGHT_GRAY, border_color=DARK_BLUE):
    """Add a detail box with title and content"""
    # Background box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    box.line.fill.background()

    # Right border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left + width - Inches(0.08), top,
        Inches(0.08), height
    )
    border.fill.solid()
    border.fill.fore_color.rgb = border_color
    border.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.1),
        width - Inches(0.4), Inches(0.3)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Content
    content_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.35),
        width - Inches(0.4), height - Inches(0.45)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = RGBColor(0x4a, 0x55, 0x68)
    p.alignment = PP_ALIGN.RIGHT

def add_event_card(slide, left, top, width, height, title, content, bg_color=DARK_BLUE):
    """Add an event card"""
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.15),
        width - Inches(0.4), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    # Content
    content_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.6),
        width - Inches(0.4), height - Inches(0.8)
    )
    tf = content_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = content
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

def create_title_slide(prs):
    """Create the title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Blue gradient background on right side
    blue_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SLIDE_WIDTH / 2, Inches(0),
        SLIDE_WIDTH / 2, SLIDE_HEIGHT
    )
    blue_bg.fill.solid()
    blue_bg.fill.fore_color.rgb = DARK_BLUE
    blue_bg.line.fill.background()

    # Arrow shape overlay
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        SLIDE_WIDTH / 2 - Inches(1), Inches(0),
        Inches(2), SLIDE_HEIGHT
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x2c, 0x52, 0x82)
    arrow.line.fill.background()

    # Green corner
    green_corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0), SLIDE_HEIGHT - Inches(1.5),
        Inches(1.5), Inches(1.5)
    )
    green_corner.fill.solid()
    green_corner.fill.fore_color.rgb = GREEN
    green_corner.line.fill.background()

    # Title: SERA 2026
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(5.5), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SERA 2026"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2),
        Inches(5.5), Inches(0.6)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "هيئة تنظيم الكهرباء"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Technical Proposal
    prop_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(5.5), Inches(0.8)
    )
    tf = prop_box.text_frame
    p = tf.paragraphs[0]
    p.text = "العرض الفني"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

def create_toc_slide(prs):
    """Create table of contents slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, 2)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.8)
    )
    set_rtl_arabic_text(title_box, "المحتوى", 44, True, DARK_BLUE)

    toc_items = [
        ("نظرة عامة على البرنامج", "03"),
        ("فعاليات الربع الأول (Q1)", "05"),
        ("فعاليات الربع الثاني (Q2)", "10"),
        ("فعاليات الربع الثالث (Q3)", "14"),
        ("فعاليات الربع الرابع (Q4)", "18"),
        ("الفعاليات الرياضية", "25"),
        ("ملخص الميزانية", "27"),
    ]

    y = Inches(1.6)
    for i, (text, page) in enumerate(toc_items):
        # Line
        line_color = GREEN if i % 2 == 0 else DARK_BLUE
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(1), y + Inches(0.5),
            Inches(11), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = line_color
        line.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(
            Inches(2), y,
            Inches(9), Inches(0.5)
        )
        set_rtl_arabic_text(text_box, text, 20, True, DARK_BLUE)

        # Page number
        page_box = slide.shapes.add_textbox(
            Inches(1), y,
            Inches(0.8), Inches(0.5)
        )
        tf = page_box.text_frame
        p = tf.paragraphs[0]
        p.text = page
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = LIGHT_BLUE
        p.alignment = PP_ALIGN.LEFT

        y += Inches(0.7)

def create_overview_slide(prs):
    """Create program overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, 3)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, "نظرة عامة على البرنامج", 32, True, DARK_BLUE)

    # Info cards
    cards = [
        ("إجمالي الفعاليات", "44 فعالية"),
        ("فعاليات Q1", "9 فعاليات"),
        ("فعاليات Q2", "8 فعاليات"),
        ("فعاليات Q3", "8 فعاليات"),
        ("فعاليات Q4", "16 فعالية"),
        ("الفعاليات الرياضية", "3 فعاليات"),
    ]

    card_width = Inches(3.5)
    card_height = Inches(1.1)
    start_x = Inches(0.8)
    y = Inches(1.5)

    for i, (title, value) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_width + Inches(0.3))
        card_y = y + row * (card_height + Inches(0.2))
        add_info_card(slide, x, card_y, card_width, card_height, title, value)

    # Detail boxes
    add_detail_box(slide, Inches(0.8), Inches(4), Inches(11.5), Inches(0.7),
                   "الموقع", "جميع الفعاليات في مدينة الرياض - المملكة العربية السعودية")
    add_detail_box(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.7),
                   "عمولة الوكالة", "15% على جميع الفعاليات")

def create_categories_slide(prs):
    """Create event categories slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, 4)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, "فئات الفعاليات", 32, True, DARK_BLUE)

    categories = [
        ("المؤتمرات والاحتفالات", "الاجتماع السنوي، احتفالات الأعياد، اليوم الوطني، حفل نهاية العام"),
        ("الاحتفالات الوطنية", "يوم التأسيس، يوم العلم السعودي، اليوم الوطني (96)"),
        ("التوعية الصحية", "التبرع بالدم، مكافحة التدخين، الصحة النفسية، السكري"),
        ("التطوير المهني", "يوم الإبداع والابتكار، برنامج تزوّد، يوم الجودة"),
        ("المسؤولية الاجتماعية", "حملة إحسان، كسوة فرح، يوم التطوع"),
        ("الفعاليات العائلية", "صيف سيرا، شتوية سيرا، يوم الطفل العالمي"),
    ]

    card_width = Inches(3.7)
    card_height = Inches(1.5)
    start_x = Inches(0.8)
    y = Inches(1.5)

    for i, (title, content) in enumerate(categories):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_width + Inches(0.2))
        card_y = y + row * (card_height + Inches(0.2))

        # Card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, card_y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.fill.background()

        # Title
        title_txt = slide.shapes.add_textbox(
            x + Inches(0.1), card_y + Inches(0.1),
            card_width - Inches(0.2), Inches(0.4)
        )
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x68, 0xd3, 0x91)
        p.alignment = PP_ALIGN.CENTER

        # Content
        content_txt = slide.shapes.add_textbox(
            x + Inches(0.1), card_y + Inches(0.5),
            card_width - Inches(0.2), card_height - Inches(0.6)
        )
        tf = content_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = content
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def create_section_divider(prs, section_title, section_subtitle, subtitle_detail, page_num):
    """Create a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, page_num)

    # Section subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5),
        Inches(11.5), Inches(0.6)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_subtitle
    p.font.size = Pt(26)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Main title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(3.2),
        Inches(11.5), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.CENTER

    # Detail
    detail_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.3),
        Inches(11.5), Inches(0.5)
    )
    tf = detail_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle_detail
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x71, 0x80, 0x96)
    p.alignment = PP_ALIGN.CENTER

def create_event_detail_slide(prs, title, date, attendance, level, venue, stage, av, services, page_num):
    """Create a detailed event slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, page_num)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, title, 32, True, DARK_BLUE)

    # Info cards
    cards = [
        ("التاريخ", date),
        ("الحضور المتوقع", attendance),
        ("المستوى", level),
    ]

    card_width = Inches(3.5)
    card_height = Inches(1)
    start_x = Inches(0.8)

    for i, (card_title, value) in enumerate(cards):
        x = start_x + i * (card_width + Inches(0.3))
        add_info_card(slide, x, Inches(1.5), card_width, card_height, card_title, value)

    # Detail boxes
    y = Inches(2.7)
    details = [
        ("المكان المقترح", venue),
        ("المسرح والديكور", stage),
        ("المتطلبات السمعية والبصرية", av),
        ("الخدمات الرئيسية", services),
    ]

    for title_txt, content in details:
        add_detail_box(slide, Inches(0.8), y, Inches(11.5), Inches(0.8), title_txt, content)
        y += Inches(0.9)

def create_two_event_slide(prs, page_title, events, page_num):
    """Create a slide with two event cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, page_num)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, page_title, 32, True, DARK_BLUE)

    card_height = Inches(2.3)
    y = Inches(1.5)

    colors = [DARK_BLUE, GREEN, RGBColor(0xb8, 0x32, 0x80), RGBColor(0xc0, 0x56, 0x21), RGBColor(0x6b, 0x46, 0xc1)]

    for i, (event_title, event_content, color_idx) in enumerate(events):
        color = colors[color_idx % len(colors)]
        add_event_card(slide, Inches(0.8), y, Inches(11.5), card_height, event_title, event_content, color)
        y += card_height + Inches(0.2)

def create_four_event_slide(prs, page_title, events, page_num, extra_detail=None):
    """Create a slide with four smaller event cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, page_num)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, page_title, 32, True, DARK_BLUE)

    card_width = Inches(5.6)
    card_height = Inches(1.4)
    start_x = Inches(0.8)
    start_y = Inches(1.5)

    for i, (event_title, event_content) in enumerate(events[:4]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_width + Inches(0.2))
        y = start_y + row * (card_height + Inches(0.2))

        # Card
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x, y, card_width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BLUE
        card.line.fill.background()

        # Title
        title_txt = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.1),
            card_width - Inches(0.3), Inches(0.4)
        )
        tf = title_txt.text_frame
        p = tf.paragraphs[0]
        p.text = event_title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x68, 0xd3, 0x91)
        p.alignment = PP_ALIGN.RIGHT

        # Content
        content_txt = slide.shapes.add_textbox(
            x + Inches(0.15), y + Inches(0.5),
            card_width - Inches(0.3), card_height - Inches(0.6)
        )
        tf = content_txt.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = event_content
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.RIGHT

    # Extra detail box at bottom
    if extra_detail:
        add_detail_box(slide, Inches(0.8), Inches(4.5), Inches(11.5), Inches(0.8),
                       extra_detail[0], extra_detail[1])

def create_three_event_slide(prs, page_title, events, page_num):
    """Create a slide with three event cards"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, page_num)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, page_title, 32, True, DARK_BLUE)

    card_height = Inches(1.5)
    y = Inches(1.5)

    colors = [DARK_BLUE, RGBColor(0xb8, 0x32, 0x80), GREEN]

    for i, (event_title, event_content, color_idx) in enumerate(events):
        color = colors[color_idx % len(colors)]
        add_event_card(slide, Inches(0.8), y, Inches(11.5), card_height, event_title, event_content, color)
        y += card_height + Inches(0.15)

def create_sports_slide(prs):
    """Create sports events slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, 26)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, "الفعاليات الرياضية", 32, True, DARK_BLUE)

    sports = [
        ("S1. بطولة كرة القدم", "الموعد: يحدد لاحقاً (3 عطلات نهاية أسبوع) | الفئة: مسابقة رياضية | الحضور: 150 شخص | المستوى: كبير", GREEN),
        ("S2. بطولة البادل", "الموعد: يحدد لاحقاً (عطلتي نهاية أسبوع) | الفئة: مسابقة رياضية | الحضور: 20 شخص | المستوى: متوسط", LIGHT_BLUE),
        ("S3. تحدي المشي + جمعة سيرا", "الموعد: ربع سنوي + سنوي | الفئة: رياضة وعافية | الحضور: 200 شخص | المستوى: متوسط", RGBColor(0xed, 0x89, 0x36)),
    ]

    y = Inches(1.5)
    card_height = Inches(1.5)

    for title, content, color in sports:
        add_event_card(slide, Inches(0.8), y, Inches(11.5), card_height, title, content, color)
        y += card_height + Inches(0.15)

def create_budget_slide(prs):
    """Create budget overview slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_top_border(slide)
    add_bottom_border(slide, 28)

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(0.6),
        Inches(11.5), Inches(0.7)
    )
    set_rtl_arabic_text(title_box, "ملخص الميزانية", 32, True, DARK_BLUE)

    # Info cards
    cards = [
        ("فعاليات Q1", "9 فعاليات"),
        ("فعاليات Q2", "8 فعاليات"),
        ("فعاليات Q3", "8 فعاليات"),
        ("فعاليات Q4", "16 فعالية"),
        ("الفعاليات الرياضية", "3 فعاليات"),
        ("إجمالي الفعاليات", "44 فعالية"),
    ]

    card_width = Inches(3.5)
    card_height = Inches(1.1)
    start_x = Inches(0.8)
    y = Inches(1.5)

    for i, (title, value) in enumerate(cards):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_width + Inches(0.3))
        card_y = y + row * (card_height + Inches(0.2))
        bg_color = GREEN if i == 5 else DARK_BLUE
        add_info_card(slide, x, card_y, card_width, card_height, title, value, bg_color)

    # Budget boxes
    add_detail_box(slide, Inches(0.8), Inches(4), Inches(11.5), Inches(0.7),
                   "الإجمالي قبل الضريبة (ريال سعودي)", "[يُحدد لاحقاً]")
    add_detail_box(slide, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.7),
                   "ضريبة القيمة المضافة (15%)", "[يُحدد لاحقاً]")

    # Total box (dark)
    total_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.7)
    )
    total_box.fill.solid()
    total_box.fill.fore_color.rgb = DARK_BLUE
    total_box.line.fill.background()

    # Border
    border = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(12.22), Inches(5.6),
        Inches(0.08), Inches(0.7)
    )
    border.fill.solid()
    border.fill.fore_color.rgb = GREEN
    border.line.fill.background()

    # Title
    title_txt = slide.shapes.add_textbox(
        Inches(1), Inches(5.65),
        Inches(11), Inches(0.3)
    )
    tf = title_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "★ الإجمالي الكلي شامل الضريبة"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

    # Value
    value_txt = slide.shapes.add_textbox(
        Inches(1), Inches(5.95),
        Inches(11), Inches(0.3)
    )
    tf = value_txt.text_frame
    p = tf.paragraphs[0]
    p.text = "[يُحدد لاحقاً]"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x68, 0xd3, 0x91)
    p.alignment = PP_ALIGN.RIGHT

def create_thank_you_slide(prs):
    """Create thank you slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Blue gradient background on right side
    blue_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        SLIDE_WIDTH / 2, Inches(0),
        SLIDE_WIDTH / 2, SLIDE_HEIGHT
    )
    blue_bg.fill.solid()
    blue_bg.fill.fore_color.rgb = DARK_BLUE
    blue_bg.line.fill.background()

    # Arrow shape overlay
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        SLIDE_WIDTH / 2 - Inches(1), Inches(0),
        Inches(2), SLIDE_HEIGHT
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = RGBColor(0x2c, 0x52, 0x82)
    arrow.line.fill.background()

    # Green corner
    green_corner = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(0), SLIDE_HEIGHT - Inches(1.5),
        Inches(1.5), Inches(1.5)
    )
    green_corner.fill.solid()
    green_corner.fill.fore_color.rgb = GREEN
    green_corner.line.fill.background()

    # Title: SERA 2026
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(5.5), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "SERA 2026"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Subtitle
    sub_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.2),
        Inches(5.5), Inches(0.6)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "هيئة تنظيم الكهرباء"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.RIGHT

    # Thank you
    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(4),
        Inches(5.5), Inches(0.8)
    )
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "شكراً لكم"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.RIGHT

def main():
    """Main function to create the presentation"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Slide 1: Title
    create_title_slide(prs)

    # Slide 2: Table of Contents
    create_toc_slide(prs)

    # Slide 3: Program Overview
    create_overview_slide(prs)

    # Slide 4: Event Categories
    create_categories_slide(prs)

    # Slide 5: Q1 Section Divider
    create_section_divider(prs, "فعاليات الربع الأول", "القسم الأول", "9 فعاليات | يناير - مارس 2026", 5)

    # Slide 6: Q1 Event 1 - Annual Meeting
    create_event_detail_slide(prs,
        "1. الاجتماع السنوي 2026",
        "3 فبراير", "300 شخص", "كبير (Major)",
        "قاعة فندق 5 نجوم (ريتز كارلتون / فور سيزونز)",
        "مسرح 8م×5م، خلفية تحمل العلامة التجارية، منصة مع شعار SERA",
        "شاشة LED P2.5 (5م×3م)، نظام صوت Line-array، 6 ميكروفونات لاسلكية، بث مباشر",
        "ترجمة فورية عربي/إنجليزي، مقدم برامج محترف، نظام تسجيل QR، تصوير (2)، فيديو (2)، مطبوعات",
        6)

    # Slide 7: Q1 Events 2-3
    create_two_event_slide(prs, "فعاليات الربع الأول", [
        ("2. يوم التأسيس - 22 فبراير",
         "المكان: مقر SERA - داخلي | الحضور: 300 شخص | الفئة: احتفال وطني\nالديكور: عروض تراثية، ثيم يوم التأسيس الوطني\nالخدمات: عروض تراثية، هدايا تذكارية، تغطية مباشرة على وسائل التواصل الاجتماعي", 0),
        ("3. يوم تقدير الموظف - 2 مارس",
         "المكان: فندق 5 نجوم - حفل عشاء | الحضور: 300 شخص | الفئة: حفل تكريم\nالديكور: مسرح للحفل، ديكور فاخر، إضاءة علوية، منصات الجوائز\nالخدمات: 12 كأس كريستال، فرقة موسيقية حية (2.5 ساعة)، صناديق هدايا، مقدم برامج محترف", 1),
    ], 7)

    # Slide 8: Q1 Events 4-5
    create_two_event_slide(prs, "فعاليات الربع الأول", [
        ("4. اليوم العالمي للمرأة - 8 مارس",
         "المكان: مكان أنيق أو قسم السيدات في فندق | الحضور: 150 شخص\nالديكور: ديكور أنثوي أنيق، تنسيقات زهور فاخرة\nالخدمات: متحدثة رئيسية + جلسة نقاشية (3 متحدثات)، ركن تصوير، حقائب هدايا، مصورة", 2),
        ("5. يوم العلم السعودي - 11 مارس",
         "المكان: مقر SERA - البهو والمنطقة الخارجية | الحضور: 300 شخص\nالديكور: علم سعودي كبير (5م+)، ثيم أخضر وأبيض وطني\nالخدمات: مراسم رفع العلم، هدايا تذكارية، تغطية مباشرة على وسائل التواصل", 1),
    ], 8)

    # Slide 9: Q1 Events 6-7
    create_two_event_slide(prs, "فعاليات الربع الأول", [
        ("6. إفطار رمضان - منتصف مارس",
         "المكان: خيمة رمضانية في فندق 5 نجوم | الحضور: 300 شخص\nالديكور: ديكور رمضاني تقليدي: فوانيس، أهلّة، مجلس VIP\nالخدمات: بوفيه إفطار فاخر، عازف عود، صناديق هدايا رمضانية، ترتيب مصلى", 0),
        ("7. يوم الأم - 21 مارس",
         "المكان: قسم خاص في مطعم أنيق | الحضور: 130 شخص\nالديكور: ثيم الزهور، منطقة تصوير\nالخدمات: هدايا زهور، هدايا خاصة للأمهات، ترفيه، تصوير فوتوغرافي", 2),
    ], 9)

    # Slide 10: Q1 Events 8-9
    create_two_event_slide(prs, "فعاليات الربع الأول", [
        ("8. المبادرة الخضراء السعودية - 27 مارس",
         "المكان: مقر SERA - منطقة المعارض | الحضور: 250 شخص\nالديكور: أكشاك بثيم الاستدامة الخضراء، عروض بيئية\nالخدمات: خبير بيئي، معارض خضراء، هدايا صديقة للبيئة، عروض استدامة", 1),
        ("9. حملة إحسان الخيرية - مارس",
         "المكان: مقر SERA - منطقة البهو | الحضور: 200 شخص\nالديكور: عروض منصة إحسان، محطات التبرع\nالخدمات: كشك توعية بالحملة، محطات تبرع QR، عرض قصص التأثير", 0),
    ], 10)

    # Slide 11: Q2 Section Divider
    create_section_divider(prs, "فعاليات الربع الثاني", "القسم الثاني", "8 فعاليات | أبريل - يونيو 2026", 11)

    # Slide 12: Q2 Event 10 - Eid Al Fitr
    create_event_detail_slide(prs,
        "10. احتفال عيد الفطر",
        "أوائل أبريل", "300 شخص", "كبير (Major)",
        "القاعة الكبرى في فندق 5 نجوم",
        "ديكور عيد احتفالي ذهبي وأبيض، مسرح كبير، قوس مدخل",
        "إنتاج كامل: شاشة LED، صوت حفلات، تصميم إضاءة",
        "فرقة عرضة + موسيقى حية، هدايا العيد، ركن للأطفال (اختياري)، مقدم برامج محترف",
        12)

    # Slide 13: Q2 Events 11-12
    create_two_event_slide(prs, "فعاليات الربع الثاني", [
        ("11. يوم الإبداع والابتكار - 21 أبريل",
         "المكان: مركز مؤتمرات | الحضور: 200 شخص\nالديكور: مسرح رئيسي + 10 أكشاك ابتكار للأقسام\nالخدمات: متحدث رئيسي، مسابقة ابتكار، 7 جوائز، تصويت تفاعلي", 3),
        ("12. اليوم العالمي للشاي - 21 مايو",
         "المكان: مقر SERA - المناطق المشتركة | الحضور: 250 شخص\nالديكور: ديكور بثيم الشاي، 3 محطات خدمة\nالخدمات: محطات شاي عربي/إنجليزي/آسيوي، مرافقات ذواقة، عروض ثقافة الشاي", 0),
    ], 13)

    # Slide 14: Q2 Events 13-17
    create_four_event_slide(prs, "فعاليات الربع الثاني", [
        ("13. كسوة فرح - مايو", "المكان: مقر SERA | الحضور: 200\nنقاط جمع، تنسيق المتطوعين، لوجستيات الفرز والتوزيع"),
        ("14. مكافحة التدخين - 31 مايو", "المكان: مقر SERA | الحضور: 200\nمتحدث صحي، محطة فحص CO، مواد توعوية"),
        ("15. يوم التبرع بالدم - 14 يونيو", "المكان: مقر SERA | الحضور: 120\nشراكة الهلال الأحمر، هدايا وشهادات للمتبرعين"),
        ("16. يوم الأب - 15 يونيو", "المكان: مطعم / قاعة | الحضور: 160\nهدايا للآباء، أنشطة بناء فريق"),
    ], 14, ("17. احتفال عيد الأضحى - منتصف يونيو", "القاعة الكبرى في فندق 5 نجوم | 300 شخص | عرضة + موسيقى حية، وليمة، هدايا العيد، مقدم برامج"))

    # Slide 15: Q3 Section Divider
    create_section_divider(prs, "فعاليات الربع الثالث", "القسم الثالث", "8 فعاليات | يوليو - سبتمبر 2026", 15)

    # Slide 16: Q3 Event 18 - SERA Summer
    create_event_detail_slide(prs,
        "18. صيف سيرا",
        "20 يوليو", "500 شخص", "كبير (Major)",
        "مكان ترفيهي - مساءً فقط (حرارة 45 درجة مئوية)",
        "مناطق متعددة: مسرح رئيسي، منطقة أطفال، منطقة طعام، ألعاب",
        "صوت وإضاءة خارجية بمستوى حفلات موسيقية",
        "ترفيه رئيسي، منطقة أطفال، 10 أكشاك ألعاب، مهرجان طعام، سحب، هدايا عائلية",
        16)

    # Slide 17: Q3 Events 19-22
    create_four_event_slide(prs, "فعاليات الربع الثالث", [
        ("19. شتوية سيرا - 19 أغسطس", "المكان: مكان داخلي مكيف | الحضور: 300\nثيم أرض العجائب الشتوية، مؤثرات ثلجية، برنامج ترفيهي"),
        ("20. الإسعافات الأولية - 13 سبتمبر", "المكان: مقر SERA | الحضور: 200\nورشة CPR، حقائب إسعافات، خبير مسعف، شهادات"),
        ("21. تزوّد - سبتمبر", "المكان: مركز تدريب | الحضور: 200\nمدربين خبراء (جلستين)، مواد تدريبية، شهادات"),
        ("22. الزهايمر - 21 سبتمبر", "المكان: مقر SERA | الحضور: 200\nمتحدث رعاية صحية، شرائط بنفسجية، عروض توعوية"),
    ], 17)

    # Slide 18: Q3 Events 23-25
    create_three_event_slide(prs, "فعاليات الربع الثالث", [
        ("23. وش دورنا - سبتمبر",
         "المكان: مقر SERA - قاعة اجتماعات | الحضور: 200 شخص\nالخدمات: ميسر محترف، تمارين تفاعلية، أدلة الأدوار، أنشطة فريق", 0),
        ("24. اليوم الوطني السعودي (96) - 23 سبتمبر",
         "المكان: مكان فاخر مع خيار خارجي | الحضور: 300 شخص\nالديكور: ديكور وطني أخضر فاخر، أعلام، رموز وطنية\nالخدمات: فرقة عرضة، عروض ثقافية، هدايا وطنية، إضاءة خضراء. يجب الحجز قبل 3-4 أشهر", 2),
        ("25. أسبوع البيئة - Q3",
         "المكان: مقر SERA - منطقة معارض (5 أيام) | الحضور: 250 شخص\nالخدمات: برنامج 5 أيام، متحدثون خبراء، زراعة أشجار، ورش بيئية، أنشطة خضراء", 0),
    ], 18)

    # Slide 19: Q4 Section Divider
    create_section_divider(prs, "فعاليات الربع الرابع", "القسم الرابع", "16 فعالية | أكتوبر - ديسمبر 2026", 19)

    # Slide 20: Q4 Events 26-28
    create_three_event_slide(prs, "فعاليات الربع الرابع", [
        ("26. اليوم العالمي للقهوة - 1 أكتوبر",
         "المكان: مقر SERA - المناطق المشتركة | الحضور: 250 شخص\nالخدمات: باريستا محترف، مراسم الدلة العربية، عروض ثقافة القهوة", 0),
        ("27. التوعية بسرطان الثدي - أكتوبر",
         "المكان: مقر SERA - منطقة الصحة | الحضور: 200 شخص\nالخدمات: خبير رعاية صحية، معلومات فحص، عناصر توعية وردية", 1),
        ("28. معرض الأمن السيبراني - أكتوبر",
         "المكان: مقر SERA / مركز مؤتمرات | الحضور: 250 شخص\nالخدمات: 6 عروض تفاعلية، 2 متحدثين خبراء، محاكاة تصيد، هدايا ترويجية", 0),
    ], 20)

    # Slide 21: Q4 Events 29-32
    create_four_event_slide(prs, "فعاليات الربع الرابع", [
        ("29. الصحة النفسية - 10 أكتوبر", "المكان: مقر SERA | الحضور: 200\nمتحدث نفسي، ورشة إدارة الضغط، منطقة استرخاء"),
        ("30. يوم الادخار - 31 أكتوبر", "المكان: مقر SERA | الحضور: 200\nخبير مالي، أدلة ادخار، مخططات ميزانية"),
        ("31. يوم الجودة - 10 نوفمبر", "المكان: فندق أعمال | الحضور: 180\n7 جوائز جودة، عروض أفضل الممارسات، متحدث خبير"),
        ("32. التطعيم ضد الإنفلونزا - نوفمبر", "المكان: مقر SERA | الحضور: 250\nطاقم طبي، لوجستيات التطعيم، هدايا للمشاركين"),
    ], 21)

    # Slide 22: Q4 Events 33-36
    create_four_event_slide(prs, "فعاليات الربع الرابع", [
        ("33. يوم السكري - 14 نوفمبر", "المكان: مقر SERA | الحضور: 200\nأخصائي غدد صماء، فحص سكر الدم، عناصر توعية زرقاء"),
        ("34. يوم الرجل - 19 نوفمبر", "المكان: قاعة فعاليات | الحضور: 200\nمسابقات فريق، هدايا تقدير"),
        ("35. يوم الطفل - 20 نوفمبر", "المكان: مكان ترفيهي عائلي | الحضور: 200\nمنشطين أطفال محترفين، محطات ألعاب، هدايا"),
        ("36. يوم التطوع - 5 ديسمبر", "المكان: موقع مجتمعي | الحضور: 80\nحافلتين نقل، قمصان متطوعين، صناديق غداء"),
    ], 22)

    # Slide 23: Q4 Events 37-40
    create_four_event_slide(prs, "فعاليات الربع الرابع", [
        ("37. مكافحة الفساد - 9 ديسمبر", "المكان: مقر SERA | الحضور: 200\nمتحدث أخلاقيات/نزاهة، توقيع تعهدات، كتيبات"),
        ("38. أنشطة التحول - Q4", "المكان: مقر SERA | الحضور: 200\nجلسات خبراء، ورش تحول، أنشطة تفكير تصميمي"),
        ("39. المكتب المثالي - Q4", "المكان: جميع مكاتب SERA | الحضور: 200\nبرنامج تقييم، لجنة تحكيم، حفل جوائز، كؤوس"),
        ("40. يوم اللغة العربية - 18 ديسمبر", "المكان: مكان ثقافي | الحضور: 200\nشاعر ضيف، ورشة خط عربي، أمسية شعرية"),
    ], 23)

    # Slide 24: Q4 Event 41 - Year End Party
    create_event_detail_slide(prs,
        "41. حفل نهاية العام",
        "أواخر ديسمبر", "300 شخص", "كبير (Major)",
        "فندق 5 نجوم - حفل فاخر",
        "مسرح فاخر، تصميم إضاءة LED، طاولات فاخرة",
        "حفل كامل: شاشة LED، صوت حفلات موسيقية، إضاءة مصممة",
        "فرقة موسيقية فاخرة + عروض، 15 جائزة سنوية، فيديو مراجعة العام، هدايا فاخرة، مقدم برامج. يجب الحجز قبل 3 أشهر",
        24)

    # Slide 25: Sports Section Divider
    create_section_divider(prs, "الفعاليات الرياضية", "القسم الخامس", "3 فعاليات رياضية", 25)

    # Slide 26: Sports Events
    create_sports_slide(prs)

    # Slide 27: Budget Section Divider
    create_section_divider(prs, "ملخص الميزانية", "القسم السادس", "نظرة عامة على التكاليف", 27)

    # Slide 28: Budget Overview
    create_budget_slide(prs)

    # Slide 29: Thank You
    create_thank_you_slide(prs)

    # Save the presentation
    output_path = '/home/user/presentations/SERA_2026_Technical_Proposal.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()
